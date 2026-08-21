#!/usr/bin/env python
"""Validate PPTX files against the pptx-meiryo skill's structural rules."""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


ALLOWED_FONTS = {"Meiryo", "Meiryo UI"}
EXPECTED_WIDTH = Inches(13.33)
EXPECTED_HEIGHT = Inches(7.5)
MIN_FONT_SIZE = Pt(10)


def emu_to_inches(value: int) -> float:
    return value / 914400


def close_enough(actual: int, expected: int, tolerance_inches: float = 0.03) -> bool:
    return abs(emu_to_inches(actual - expected)) <= tolerance_inches


def validate(path: Path) -> list[str]:
    prs = Presentation(path)
    errors: list[str] = []

    if not close_enough(prs.slide_width, EXPECTED_WIDTH):
        errors.append(
            f"Deck: width is {emu_to_inches(prs.slide_width):.2f}in, expected 13.33in"
        )
    if not close_enough(prs.slide_height, EXPECTED_HEIGHT):
        errors.append(
            f"Deck: height is {emu_to_inches(prs.slide_height):.2f}in, expected 7.50in"
        )

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape_idx, shape in enumerate(slide.shapes, 1):
            label = f"Slide {slide_idx}, shape {shape_idx}"

            if shape.left < 0:
                errors.append(f"{label}: left edge is outside slide")
            if shape.top < 0:
                errors.append(f"{label}: top edge is outside slide")
            if shape.left + shape.width > prs.slide_width:
                errors.append(f"{label}: right edge exceeds slide width")
            if shape.top + shape.height > prs.slide_height:
                errors.append(f"{label}: bottom edge exceeds slide height")

            if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
                continue

            text_frame = shape.text_frame
            if getattr(text_frame, "word_wrap", None) is False:
                errors.append(f"{label}: word_wrap is disabled")

            for para_idx, paragraph in enumerate(text_frame.paragraphs, 1):
                if paragraph.level > 1:
                    errors.append(f"{label}, paragraph {para_idx}: bullet level exceeds 2")

                for run_idx, run in enumerate(paragraph.runs, 1):
                    run_label = f"{label}, paragraph {para_idx}, run {run_idx}"
                    if run.text.strip() and run.font.name not in ALLOWED_FONTS:
                        errors.append(f"{run_label}: font is not Meiryo")
                    if run.font.size is not None and run.font.size < MIN_FONT_SIZE:
                        errors.append(f"{run_label}: font size is below 10pt")

    return errors


def main() -> int:
    parser = argparse.ArgumentParser(description="Validate pptx-meiryo PPTX rules.")
    parser.add_argument("pptx", type=Path, help="Path to a .pptx file")
    parser.add_argument(
        "--max-errors",
        type=int,
        default=50,
        help="Maximum number of validation errors to print. Use 0 for all.",
    )
    args = parser.parse_args()

    if not args.pptx.exists():
        print(f"File not found: {args.pptx}", file=sys.stderr)
        return 2

    errors = validate(args.pptx)
    if errors:
        print("Validation failed:")
        shown_errors = errors if args.max_errors == 0 else errors[: args.max_errors]
        for error in shown_errors:
            print(f"- {error}")
        remaining = len(errors) - len(shown_errors)
        if remaining > 0:
            print(f"... {remaining} more issue(s) omitted. Re-run with --max-errors 0 to show all.")
        return 1

    print("Validation passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
